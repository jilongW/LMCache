from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

OUT_DIR = Path("/home/user/jlwang/models/lmcache/kvweave/lmcache-main/artifacts/cache_plots")
OUT_DIR.mkdir(parents=True, exist_ok=True)
OUT_FILE = OUT_DIR / "cache_quant_story_editable.pptx"

DATA_2Q = {
    "NoQuant": [
        {"user": "User1", "req": "req1", "ttft": 74700.94, "load": 0, "store": 40960},
        {"user": "User2", "req": "req1", "ttft": 31473.69, "load": 0, "store": 40960},
        {"user": "User1", "req": "req2", "ttft": 73884.53, "load": 0, "store": 40960},
        {"user": "User2", "req": "req2", "ttft": 30771.74, "load": 0, "store": 40960},
    ],
    "Quant": [
        {"user": "User1", "req": "req1", "ttft": 40689.17, "load": 0, "store": 40960},
        {"user": "User2", "req": "req1", "ttft": 92122.75, "load": 0, "store": 40960},
        {"user": "User1", "req": "req2", "ttft": 22433.72, "load": 40960, "store": 0},
        {"user": "User2", "req": "req2", "ttft": 22433.72, "load": 40960, "store": 0},
    ],
}

DATA_3Q = {
    "NoQuant": [
        {"req": "req1", "input": 40960, "output": 128, "ttft": 31306.60, "tpot": 98.44, "load": 0, "l2_load": 256, "store": 40960},
        {"req": "req2", "input": 40960, "output": 128, "ttft": 30869.98, "tpot": 103.67, "load": 0, "l2_load": 0, "store": 40960},
        {"req": "req3", "input": 40960, "output": 128, "ttft": 30351.83, "tpot": 103.64, "load": 0, "l2_load": 0, "store": 40960},
        {"req": "req1", "input": 40960, "output": 128, "ttft": 30552.87, "tpot": 103.65, "load": 256, "l2_load": 0, "store": 40704},
        {"req": "req2", "input": 40960, "output": 128, "ttft": 30619.58, "tpot": 103.77, "load": 256, "l2_load": 256, "store": 40448},
        {"req": "req3", "input": 40960, "output": 128, "ttft": 30746.89, "tpot": 99.22, "load": 0, "l2_load": 0, "store": 40960},
    ],
    "Quant": [
        {"req": "req1", "input": 40960, "output": 128, "ttft": 40523.25, "tpot": 99.16, "load": 0, "l2_load": 0, "store": 40960},
        {"req": "req2", "input": 40960, "output": 128, "ttft": 41785.93, "tpot": 99.11, "load": 0, "l2_load": 0, "store": 40960},
        {"req": "req3", "input": 40960, "output": 128, "ttft": 40347.48, "tpot": 99.03, "load": 0, "l2_load": 0, "store": 40960},
        {"req": "req1", "input": 40960, "output": 128, "ttft": 12990.40, "tpot": 98.17, "load": 10240, "l2_load": 30720, "store": 0},
        {"req": "req2", "input": 40960, "output": 128, "ttft": 12265.84, "tpot": 99.21, "load": 4096, "l2_load": 36864, "store": 0},
        {"req": "req3", "input": 40960, "output": 128, "ttft": 12332.00, "tpot": 98.45, "load": 40960, "l2_load": 0, "store": 0},
    ],
}


def add_title(slide, title, subtitle=""):
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.6))
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(30)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(33, 33, 33)

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.5))
        tf2 = sub.text_frame
        tf2.text = subtitle
        tf2.paragraphs[0].font.size = Pt(14)
        tf2.paragraphs[0].font.color.rgb = RGBColor(90, 90, 90)


def retention(items):
    vals = []
    for item in items:
        total = item["load"] + item["store"]
        vals.append(0 if total == 0 else item["load"] / total)
    return sum(vals) / len(vals)


def add_summary_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "LMCache Quant Reuse Summary", "All elements are editable in PowerPoint")

    bullets = [
        "Without Quant: repeated requests still rebuild cache (load ~ 0, store high)",
        "With Quant: repeated requests can reuse cache (load rises, TTFT drops)",
        "2Q repeated round: Quant reaches full reuse (load=40960, store=0)",
        "3Q repeated round: Quant has partial reuse (load=9216, store=31744)",
    ]

    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(11.8), Inches(2.6))
    tf = box.text_frame
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(20)
        p.level = 0

    card_w = Inches(5.7)
    card_h = Inches(2.3)

    def add_card(x, y, title, n_val, q_val):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(244, 246, 247)
        shape.line.color.rgb = RGBColor(210, 214, 220)

        t = shape.text_frame
        t.text = title
        t.paragraphs[0].font.size = Pt(16)
        t.paragraphs[0].font.bold = True
        p1 = t.add_paragraph()
        p1.text = f"NoQuant retention: {n_val:.1%}"
        p1.font.size = Pt(14)
        p2 = t.add_paragraph()
        p2.text = f"Quant retention:   {q_val:.1%}"
        p2.font.size = Pt(14)

    two_no = retention(DATA_2Q["NoQuant"][2:])
    two_q = retention(DATA_2Q["Quant"][2:])
    three_no = retention(DATA_3Q["NoQuant"][3:])
    three_q = retention(DATA_3Q["Quant"][3:])

    add_card(Inches(0.7), Inches(4.2), "2Q repeated round", two_no, two_q)
    add_card(Inches(6.5), Inches(4.2), "3Q repeated round", three_no, three_q)


def add_sequence_table_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "3Q Sequential Timeline (Editable Table)", "One user sends req1 -> req2 -> req3 in sequence")

    rows = 3
    cols = 13
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.4), Inches(1.5), Inches(12.5), Inches(2.2))
    table = table_shape.table

    headers = ["Mode"]
    for index in range(1, 7):
        headers.extend([f"req{index}", "TTFT(ms)"])
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(230, 235, 240)

    def fill_row(r, mode, seq):
        table.cell(r, 0).text = mode
        for i, item in enumerate(seq):
            table.cell(r, 1 + i * 2).text = item["req"]
            table.cell(r, 2 + i * 2).text = f"{item['ttft']:.2f}"

    fill_row(1, "NoQuant", DATA_3Q["NoQuant"])
    fill_row(2, "Quant", DATA_3Q["Quant"])

    for r in range(rows):
        for c in range(cols):
            for p in table.cell(r, c).text_frame.paragraphs:
                p.font.size = Pt(11 if r == 0 else 10)
                if r == 0:
                    p.font.bold = True

    # Editable timeline blocks below the table.
    # Each mode uses three rows grouped by request id (req1/req2/req3).
    base_y = Inches(4.0)
    line_gap = Inches(1.55)
    # Scale by total sequence duration so each lane fits in one slide page.
    seq_totals = [sum(item["ttft"] for item in mode) for mode in DATA_3Q.values()]
    max_total_ttft = max(seq_totals)
    timeline_start = 1.95
    timeline_width = 9.8
    scale = timeline_width / max_total_ttft

    # Light backdrop to make the timeline region visually grouped.
    timeline_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.38),
        Inches(3.82),
        Inches(12.45),
        Inches(3.20),
    )
    timeline_bg.fill.solid()
    timeline_bg.fill.fore_color.rgb = RGBColor(246, 248, 251)
    timeline_bg.fill.transparency = 0.0
    timeline_bg.line.color.rgb = RGBColor(222, 227, 234)
    timeline_bg.line.width = Pt(1.0)

    # Section header bar inside timeline panel.
    hdr = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.55),
        Inches(3.95),
        Inches(12.1),
        Inches(0.34),
    )
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = RGBColor(233, 239, 247)
    hdr.line.color.rgb = RGBColor(210, 221, 236)
    hdr.text_frame.text = "Sequential Timeline Area (editable)"
    hdr.text_frame.paragraphs[0].font.size = Pt(10)
    hdr.text_frame.paragraphs[0].font.bold = True
    hdr.text_frame.paragraphs[0].font.color.rgb = RGBColor(50, 70, 95)

    # Lightweight legend chips.
    chip1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(10.05),
        Inches(3.99),
        Inches(1.2),
        Inches(0.26),
    )
    chip1.fill.solid()
    chip1.fill.fore_color.rgb = RGBColor(95, 108, 122)
    chip1.line.color.rgb = RGBColor(95, 108, 122)
    chip1.text_frame.text = "NoQuant"
    chip1.text_frame.paragraphs[0].font.size = Pt(8)
    chip1.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    chip2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(11.35),
        Inches(3.99),
        Inches(1.2),
        Inches(0.26),
    )
    chip2.fill.solid()
    chip2.fill.fore_color.rgb = RGBColor(29, 78, 137)
    chip2.line.color.rgb = RGBColor(29, 78, 137)
    chip2.text_frame.text = "Quant"
    chip2.text_frame.paragraphs[0].font.size = Pt(8)
    chip2.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    def add_lane(mode_name, items, y, color):
        req_rows = {
            "req1": y,
            "req2": y + Inches(0.42),
            "req3": y + Inches(0.84),
        }

        for req_name, req_y in req_rows.items():
            row_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1.86),
                req_y - Inches(0.02),
                Inches(10.05),
                Inches(0.34),
            )
            row_bg.fill.solid()
            row_bg.fill.fore_color.rgb = RGBColor(252, 253, 255)
            row_bg.line.color.rgb = RGBColor(234, 238, 244)
            row_bg.line.width = Pt(0.7)

            req_label = slide.shapes.add_textbox(Inches(1.32), req_y + Inches(0.08), Inches(0.5), Inches(0.2))
            req_label.text_frame.text = req_name
            req_label.text_frame.paragraphs[0].font.size = Pt(8)

        title = slide.shapes.add_textbox(Inches(0.5), y + Inches(0.18), Inches(1.1), Inches(0.3))
        title.text_frame.text = mode_name
        title.text_frame.paragraphs[0].font.size = Pt(11)
        title.text_frame.paragraphs[0].font.bold = True

        x = timeline_start
        req_counts = {"req1": 0, "req2": 0, "req3": 0}
        first_mid = {}
        for item in items:
            req_name = item["req"]
            req_counts[req_name] += 1
            w = item["ttft"] * scale
            block_y = req_rows[req_name]
            block = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x),
                block_y,
                Inches(w),
                Inches(0.30),
            )
            block.fill.solid()
            block.fill.fore_color.rgb = color
            block.line.color.rgb = RGBColor(255, 255, 255)
            block.text_frame.text = f"run{req_counts[req_name]} | {item['ttft']/1000:.2f}s"
            block.text_frame.paragraphs[0].font.size = Pt(8)
            block.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

            mid = x + w / 2
            if req_counts[req_name] == 1:
                first_mid[req_name] = mid
            else:
                conn_y = block_y + Inches(0.34)
                arr = slide.shapes.add_connector(
                    1,
                    Inches(first_mid[req_name]),
                    conn_y,
                    Inches(mid),
                    conn_y,
                )
                arr.line.color.rgb = color
                arr.line.width = Pt(1.25)
            x += w + 0.05

        # Draw lane end marker to make fitting boundary explicit.
        end_line_x = timeline_start + timeline_width
        end_line = slide.shapes.add_connector(
            1,
            Inches(end_line_x),
            y - Inches(0.08),
            Inches(end_line_x),
            y + Inches(1.24),
        )
        end_line.line.color.rgb = RGBColor(195, 195, 195)
        end_line.line.width = Pt(0.8)

    add_lane("NoQuant", DATA_3Q["NoQuant"], base_y, RGBColor(95, 108, 122))
    add_lane("Quant", DATA_3Q["Quant"], base_y + line_gap, RGBColor(29, 78, 137))

    # Divider between mode groups.
    mid_div = slide.shapes.add_connector(
        1,
        Inches(0.7),
        base_y + Inches(1.40),
        Inches(12.55),
        base_y + Inches(1.40),
    )
    mid_div.line.color.rgb = RGBColor(210, 217, 226)
    mid_div.line.width = Pt(1.0)


def add_raw_table_slide(prs, title, dataset, start_y=1.4):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, title, "Directly editable numbers for copy/paste")

    rows = 1 + sum(len(v) for v in dataset.values())
    cols = 9
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(start_y), Inches(11.8), Inches(4.7))
    table = table_shape.table

    headers = [
        "Mode", "Request", "Input", "Output", "TTFT(ms)", "TPOT(ms)",
        "L1 load tokens", "L2 load tokens", "store tokens"
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(230, 235, 240)

    r = 1
    for mode_name, items in dataset.items():
        for item in items:
            table.cell(r, 0).text = mode_name
            table.cell(r, 1).text = item["req"]
            table.cell(r, 2).text = str(item.get("input", 40960))
            table.cell(r, 3).text = str(item.get("output", 128))
            table.cell(r, 4).text = f"{item['ttft']:.2f}"
            table.cell(r, 5).text = f"{item.get('tpot', 0):.2f}"
            table.cell(r, 6).text = str(item["load"])
            table.cell(r, 7).text = str(item.get("l2_load", 0))
            table.cell(r, 8).text = str(item["store"])
            r += 1

    for rr in range(rows):
        for cc in range(cols):
            for p in table.cell(rr, cc).text_frame.paragraphs:
                p.font.size = Pt(12 if rr == 0 else 11)
                if rr == 0:
                    p.font.bold = True


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_summary_slide(prs)
    add_sequence_table_slide(prs)
    add_raw_table_slide(prs, "2Q Raw Data", DATA_2Q)
    add_raw_table_slide(prs, "3Q Raw Data", DATA_3Q)

    prs.save(OUT_FILE)
    print(f"generated: {OUT_FILE}")


if __name__ == "__main__":
    main()
